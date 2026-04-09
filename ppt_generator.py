"""
College Management System - Project Presentation Generator
Generates a comprehensive PowerPoint covering all project documentation topics.
"""
import io
import os
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Diagram Paths ─────────────────────────────────────────────
_BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DIAGRAMS_DIR = os.path.join(_BASE_DIR, 'static', 'diagrams')

def _diagram(name):
    """Return full path to a diagram image, or None if not found."""
    path = os.path.join(DIAGRAMS_DIR, name)
    return path if os.path.exists(path) else None

# ── Theme Colors ──────────────────────────────────────────────
ROYAL_BLUE = RGBColor(65, 105, 225)
DARK_BLUE = RGBColor(11, 19, 43)
GOLD = RGBColor(255, 215, 0)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 210)
DARK_GRAY = RGBColor(40, 50, 70)
ACCENT = RGBColor(212, 175, 55)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ── Helper Functions ──────────────────────────────────────────

def _set_slide_bg(slide, r, g, b):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def _add_shape_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        shape.fill.fore_color.brightness = alpha
    return shape


def _add_text_box(slide, left, top, width, height, text, font_size=18,
                  color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def _add_bullet_slide(slide, items, left, top, width, height,
                      font_size=16, color=LIGHT_GRAY, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = spacing
        p.level = 0
    return tf


def _add_image_slide(prs, title, img_path, caption=None):
    """Add a slide with a full-page diagram image and an optional caption."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, 11, 19, 43)
    _add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
    _add_text_box(slide, Inches(0.5), Inches(0.15), Inches(12), Inches(0.55),
                  title, font_size=22, color=GOLD, bold=True)
    if img_path and os.path.exists(img_path):
        # Leave space for title (0.7") and optional caption (0.4" at bottom)
        cap_h = Inches(0.35) if caption else Inches(0)
        img_top = Inches(0.75)
        img_h = SLIDE_HEIGHT - img_top - cap_h - Inches(0.1)
        pic = slide.shapes.add_picture(img_path, Inches(0.4), img_top,
                                       width=Inches(12.5), height=img_h)
    if caption:
        _add_text_box(slide, Inches(0.5), Inches(7.1), Inches(12), Inches(0.35),
                      caption, font_size=12, color=LIGHT_GRAY,
                      alignment=PP_ALIGN.CENTER)
    return slide


def _section_title_slide(prs, number, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, 11, 19, 43)
    # Accent bar
    _add_shape_rect(slide, Inches(0), Inches(3.0), Inches(13.333), Inches(1.5), ROYAL_BLUE)
    # Section number
    _add_text_box(slide, Inches(0.8), Inches(3.1), Inches(1.5), Inches(1.3),
                  str(number), font_size=48, color=GOLD, bold=True,
                  alignment=PP_ALIGN.CENTER)
    # Section title
    _add_text_box(slide, Inches(2.5), Inches(3.15), Inches(9), Inches(1.2),
                  title, font_size=36, color=WHITE, bold=True)
    # Decorative line
    _add_shape_rect(slide, Inches(2.5), Inches(4.55), Inches(3), Inches(0.05), GOLD)
    return slide


def _content_slide(prs, title, bullets, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, 11, 19, 43)
    # Top bar
    _add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
    # Title
    _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                  title, font_size=28, color=GOLD, bold=True)
    # Underline
    _add_shape_rect(slide, Inches(0.8), Inches(1.1), Inches(2.5), Inches(0.04), ROYAL_BLUE)

    y_start = 1.5
    if subtitle:
        _add_text_box(slide, Inches(0.8), Inches(y_start), Inches(11), Inches(0.5),
                      subtitle, font_size=18, color=LIGHT_GRAY, bold=False)
        y_start += 0.6

    _add_bullet_slide(slide, bullets, Inches(1.0), Inches(y_start),
                      Inches(11), Inches(5.5 - y_start + 1))
    return slide


def _table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, 11, 19, 43)
    _add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
    _add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
                  title, font_size=26, color=GOLD, bold=True)

    num_rows = min(len(rows) + 1, 16)
    num_cols = len(headers)
    tbl_left = Inches(0.8)
    tbl_top = Inches(1.3)
    tbl_width = Inches(11.7)
    tbl_height = Inches(0.4 * num_rows)

    shape = slide.shapes.add_table(num_rows, num_cols, tbl_left, tbl_top, tbl_width, tbl_height)
    table = shape.table

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ROYAL_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.color.rgb = WHITE
            paragraph.font.bold = True
            paragraph.font.name = 'Calibri'

    for ri, row in enumerate(rows[:num_rows - 1]):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(28, 37, 65) if ri % 2 == 0 else RGBColor(20, 28, 50)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = LIGHT_GRAY
                paragraph.font.name = 'Calibri'
    return slide


# ═══════════════════════════════════════════════════════════════
#  MAIN GENERATOR
# ═══════════════════════════════════════════════════════════════

def generate_project_ppt(db_tables=None):
    """Generate the full project presentation. Returns BytesIO buffer."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # ──────────────────────────────────────────────────────────
    # SLIDE 1 — TITLE SLIDE
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, 11, 19, 43)
    _add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(7.5),
                    DARK_BLUE)
    _add_shape_rect(slide, Inches(0), Inches(2.2), Inches(13.333), Inches(3.2), ROYAL_BLUE)
    _add_text_box(slide, Inches(1), Inches(0.6), Inches(11), Inches(0.6),
                  "A PROJECT REPORT ON", font_size=18, color=GOLD, bold=False,
                  alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.0),
                  "COLLEGE MANAGEMENT SYSTEM", font_size=44, color=WHITE, bold=True,
                  alignment=PP_ALIGN.CENTER)
    _add_shape_rect(slide, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.05), GOLD)
    _add_text_box(slide, Inches(1), Inches(3.7), Inches(11), Inches(0.7),
                  "A Web-Based Solution for Academic Administration",
                  font_size=20, color=LIGHT_GRAY, bold=False, alignment=PP_ALIGN.CENTER)
    _add_text_box(slide, Inches(1), Inches(5.8), Inches(5), Inches(0.5),
                  "Submitted By: Dhruv Chauhan", font_size=16, color=GOLD, bold=True)
    _add_text_box(slide, Inches(7), Inches(5.8), Inches(5), Inches(0.5),
                  f"Date: {datetime.now().strftime('%B %d, %Y')}", font_size=16,
                  color=GOLD, bold=True, alignment=PP_ALIGN.RIGHT)
    _add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
                  "Built with Flask • PostgreSQL • Bootstrap 5 • Chart.js",
                  font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # ──────────────────────────────────────────────────────────
    # SLIDE 2 — TABLE OF CONTENTS
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, 11, 19, 43)
    _add_shape_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)
    _add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.7),
                  "TABLE OF CONTENTS", font_size=30, color=GOLD, bold=True)
    _add_shape_rect(slide, Inches(0.8), Inches(1.1), Inches(3), Inches(0.04), ROYAL_BLUE)
    toc = [
        "1.  Introduction",
        "2.  About Project — Existing & Proposed System",
        "3.  Requirement Analysis — Hardware, Software, Data Dictionary, ER Diagram",
        "4.  UML Diagrams — Use Case Diagram, Activity Diagram",
        "5.  Flow Chart",
        "6.  System Development — Screenshots of Input/Output Forms",
        "7.  Implementation",
        "8.  Bibliography",
    ]
    _add_bullet_slide(slide, toc, Inches(1.2), Inches(1.6), Inches(10), Inches(5),
                      font_size=20, color=WHITE, spacing=Pt(14))

    # ──────────────────────────────────────────────────────────
    # 1. INTRODUCTION (Slides 3-4)
    # ──────────────────────────────────────────────────────────
    _section_title_slide(prs, 1, "Introduction")

    _content_slide(prs, "1. Introduction", [
        "• The College Management System (CMS) is a comprehensive web-based application",
        "  designed to automate and streamline the daily operations of educational institutions.",
        "",
        "• It replaces traditional paper-based systems with a digital platform that manages",
        "  students, faculty, attendance, examinations, fees, timetables, and notices.",
        "",
        "• The system provides role-based dashboards for Admin, Faculty, Student, and Accountant,",
        "  ensuring each user sees only the functionalities relevant to their role.",
        "",
        "• Key Objectives:",
        "  – Eliminate manual record-keeping and reduce human errors",
        "  – Provide real-time access to academic data for all stakeholders",
        "  – Enable QR-based attendance with fraud detection",
        "  – Offer secure, scalable, and deployable-to-cloud architecture",
    ])

    _content_slide(prs, "1. Introduction — Scope & Purpose", [
        "• Scope of the Project:",
        "  – Complete student lifecycle management (admission → graduation)",
        "  – Faculty management with subject assignment and timetable scheduling",
        "  – Attendance tracking via traditional and QR-code methods",
        "  – Examination management with internal/external marks and grading",
        "  – Fee structure management with payment tracking and receipt generation",
        "  – Notice board with category-based announcements",
        "  – Comprehensive reporting and analytics dashboards",
        "",
        "• Target Users:",
        "  – College administrators, faculty members, students, and accountants",
        "",
        "• Deployment: Cloud-ready (Render / Heroku) with PostgreSQL support",
    ])

    # ──────────────────────────────────────────────────────────
    # 2. ABOUT PROJECT (Slides 5-7)
    # ──────────────────────────────────────────────────────────
    _section_title_slide(prs, 2, "About Project")

    _content_slide(prs, "2.1 Existing System", [
        "• Most colleges still rely on manual, paper-based administration:",
        "",
        "  – Student records maintained in physical registers",
        "  – Attendance taken on paper sheets, prone to proxy attendance",
        "  – Fee receipts issued manually, difficult to track payment status",
        "  – Exam results compiled in Excel spreadsheets, error-prone",
        "  – Timetable changes communicated verbally or on notice boards",
        "",
        "• Problems with the Existing System:",
        "  – Slow data retrieval and reporting",
        "  – No real-time access for students or parents",
        "  – High risk of data loss (fire, misplacement, damage)",
        "  – Duplicate entry across multiple registers",
        "  – No centralized communication channel",
    ])

    _content_slide(prs, "2.2 Proposed System", [
        "• The CMS is a fully digital, web-based solution that addresses all limitations:",
        "",
        "  – Centralized database (PostgreSQL) with real-time CRUD operations",
        "  – Auto-generated enrollment numbers ({COURSE}{YEAR}{SEQ}) for students",
        "  – Secure QR-based attendance with encrypted time-limited tokens",
        "  – Online fee payment tracking with receipt generation",
        "  – Role-based access control with session management and password hashing",
        "  – Responsive UI with modern dark theme (Royal Blue & Gold)",
        "",
        "• Advantages of the Proposed System:",
        "  – 24/7 access from any device with a browser",
        "  – Instant report generation (attendance, fees, results)",
        "  – Reduced paperwork and administrative overhead",
        "  – Fraud detection in attendance via IP/device tracking",
        "  – Scalable cloud deployment (Render with Gunicorn)",
    ])

    # ──────────────────────────────────────────────────────────
    # 3. REQUIREMENT ANALYSIS (Slides 8-14)
    # ──────────────────────────────────────────────────────────
    _section_title_slide(prs, 3, "Requirement Analysis")

    # 3.1 Hardware
    _table_slide(prs, "3.1 Hardware Requirements", ["Component", "Minimum Specification", "Recommended"],
        [
            ["Processor", "Intel Core i3 / AMD Ryzen 3", "Intel Core i5 / Ryzen 5 or higher"],
            ["RAM", "4 GB", "8 GB or more"],
            ["Hard Disk", "256 GB SSD", "512 GB SSD"],
            ["Network", "10 Mbps broadband", "50 Mbps or higher"],
            ["Display", "1366 × 768", "1920 × 1080 (Full HD)"],
            ["Server (Cloud)", "1 vCPU, 512 MB RAM", "2 vCPU, 1 GB RAM (Render)"],
        ])

    # 3.2 Software
    _table_slide(prs, "3.2 Software Requirements", ["Category", "Technology", "Version"],
        [
            ["Programming Language", "Python", "3.10+"],
            ["Web Framework", "Flask", "3.0.0"],
            ["Database (Production)", "PostgreSQL", "14+"],
            ["Database (Development)", "MySQL", "8.0"],
            ["Frontend Framework", "Bootstrap", "5.3.2"],
            ["Charting Library", "Chart.js", "4.x"],
            ["Real-time Communication", "Flask-SocketIO", "5.3.6"],
            ["PDF Generation", "ReportLab", "4.0.9"],
            ["Document Generation", "python-docx / python-pptx", "1.1.0 / 0.6.21"],
            ["QR Code Generation", "qrcode + Pillow", "7.4.2 / 10.1.0"],
            ["Web Server (Production)", "Gunicorn + Eventlet", "21.2.0 / 0.33.3"],
            ["Operating System", "Windows / Linux / macOS", "Any"],
            ["Browser", "Chrome / Firefox / Edge", "Latest"],
        ])

    # 3.3 Data Dictionary
    if db_tables:
        for tbl in db_tables[:15]:
            tname = tbl['table_name']
            cols = tbl['columns']
            rows_data = []
            for c in cols:
                rows_data.append([
                    str(c['column_name']),
                    str(c['data_type']),
                    str(c['is_nullable']),
                    str(c['column_default'])[:40] if c['column_default'] else '-',
                    str(c['character_maximum_length']) if c['character_maximum_length'] else '-'
                ])
            _table_slide(prs, f"3.3 Data Dictionary — {tname.upper()}",
                         ["Column Name", "Data Type", "Nullable", "Default", "Max Length"],
                         rows_data)
    else:
        _content_slide(prs, "3.3 Data Dictionary", [
            "• The system database contains 20+ normalized tables including:",
            "  – roles, users, departments, courses, subjects",
            "  – students, faculty, faculty_subject_assignment",
            "  – attendance, attendance_sessions, attendance_records",
            "  – marks, fee_structure, fee_payments, notices, timetable",
            "  – chat_conversations, chat_messages, chat_files",
            "  – attendance_audit_log, ip_whitelist, exams, payments",
            "",
            "• All tables have proper primary keys, foreign keys, indexes,",
            "  and constraints for data integrity.",
        ])

    # 3.4 ER Diagram — embed actual image
    _add_image_slide(prs, "3.4 Entity Relationship (ER) Diagram",
                     _diagram('er_diagram.png'),
                     "ER Diagram — College Management System Database")

    # 3.5 ER Diagram relationships (text summary fallback)
    _content_slide(prs, "3.5 ER Diagram — Key Relationships", [
        "• users.role_id → roles.id  (Each user has exactly one role)",
        "• students.course_id → courses.id  (Each student belongs to one course)",
        "• courses.department_id → departments.id  (Each course belongs to one department)",
        "• subjects.course_id → courses.id  (Subjects are per course & semester)",
        "• attendance.student_id → students.id  |  attendance.subject_id → subjects.id",
        "• marks.student_id → students.id  |  marks.subject_id → subjects.id",
        "• fee_payments.student_id → students.id",
        "• fee_payments.fee_structure_id → fee_structure.id",
        "• timetable → courses, subjects, faculty  (Scheduling relationships)",
        "• attendance_sessions → faculty, subjects, courses  (QR session links)",
        "• attendance_records.session_id → attendance_sessions.id",
        "• notices.user_id → users.id  (Posted by admin/faculty)",
    ])

    # ──────────────────────────────────────────────────────────
    # 4. UML DIAGRAMS (Slides 15-18)
    # ──────────────────────────────────────────────────────────
    _section_title_slide(prs, 4, "UML Diagrams")

    # 4.1 Use Case Diagram — actual image
    _add_image_slide(prs, "4.1 UML Use Case Diagram",
                     _diagram('use_case_diagram.png'),
                     "Use Case Diagram — Admin, Faculty, Student, Accountant interactions")

    # Use Case text detail
    _content_slide(prs, "4.1 Use Case — Actor Summary", [
        "• ADMIN: Manage Students, Faculty, Academic, Attendance, Exams, Fees, Timetable, Notices, Reports",
        "",
        "• FACULTY: Take Attendance (Manual / QR), Generate QR Session,",
        "  View Assigned Subjects, Enter Exam Marks, View Timetable, Post Notices",
        "",
        "• STUDENT: View Dashboard, Scan QR for Attendance, View Attendance Summary,",
        "  View Exam Results, Pay Fees, View Timetable, View Notices",
        "",
        "• ACCOUNTANT: Set Fee Structure, Record Payments, Generate Receipts, View Fee Reports",
    ])

    # 4.2 Class Diagram — actual image
    _add_image_slide(prs, "4.2 UML Class Diagram",
                     _diagram('class_diagram.png'),
                     "Class Diagram — System classes with attributes, methods & relationships")

    # 4.3 Sequence Diagram — actual image
    _add_image_slide(prs, "4.3 UML Sequence Diagram — Login & Authentication",
                     _diagram('sequence_diagram.png'),
                     "Sequence Diagram — Role-based login flow between Browser, Flask, and Database")

    # 4.4 Activity Diagram — actual image
    _add_image_slide(prs, "4.4 UML Activity Diagram — QR Attendance Process",
                     _diagram('activity_diagram.png'),
                     "Activity Diagram — Faculty generates QR → Student scans → System validates → Attendance recorded")

    # Activity detail text slide
    _content_slide(prs, "4.4 Activity Diagram — Process Detail", [
        "• Faculty Flow:",
        "  Start → Select Subject/Course/Semester → Set Duration",
        "  → System generates encrypted QR token → QR displayed on screen",
        "  → Countdown timer starts → Students scan",
        "  → Session expires or Faculty closes → Attendance saved → End",
        "",
        "• Student Flow:",
        "  Start → Open QR Scanner page → Scan QR code OR enter token manually",
        "  → System validates: token + expiry + duplicate check",
        "  → [Valid] → Attendance marked as 'Verified' → End",
        "  → [Invalid/Expired] → Error shown → Audit log entry → End",
        "  → [Duplicate] → Already marked message → End",
    ])

    # ──────────────────────────────────────────────────────────
    # 5. FLOW CHART (Slides 19-20)
    # ──────────────────────────────────────────────────────────
    _section_title_slide(prs, 5, "Flow Chart")

    # 5. System Flowchart — actual image
    _add_image_slide(prs, "5. System Flow Chart",
                     _diagram('system_flowchart.png'),
                     "System Flowchart — Login → Role Check → Feature Access → Logout")

    _content_slide(prs, "5. Flow Chart — Module Interactions", [
        "• Admin Dashboard → Students / Faculty / Academic / Attendance / Exams / Fees / Notices",
        "",
        "• Student Module: Registration → Profile → Attendance → Results → Fees → Timetable",
        "",
        "• Faculty Module: Profile → Subject Assignment → Take Attendance → Enter Marks → Timetable",
        "",
        "• Attendance Module: Manual Entry ──OR── QR Generation → QR Scan → Validation → Record",
        "",
        "• Fees Module: Structure Setup → Student View → Payment → Receipt",
        "",
        "• Reports Module: Select Type → Filter Parameters → Generate → Export (CSV/PDF)",
    ])

    # ──────────────────────────────────────────────────────────
    # 6. SYSTEM DEVELOPMENT — Screenshots (Slides 21-24)
    # ──────────────────────────────────────────────────────────
    _section_title_slide(prs, 6, "System Development")

    _content_slide(prs, "6.1 Input Forms", [
        "• Login Form: Username/Email + Password input with role-based redirection",
        "• Student Registration Form: First Name, Last Name, Email, Phone, DOB,",
        "  Gender, Address, Course (dropdown), Semester, Photo Upload",
        "• Faculty Registration Form: Employee ID, Name, Email, Department, Designation",
        "• Attendance Form: Select Course → Subject → Date → Mark Present/Absent/Leave",
        "• QR Attendance Form: Select Subject, Course, Semester, Duration (5-30 min)",
        "• Exam Marks Entry: Select Course/Subject/Semester → Enter Internal & External marks",
        "• Fee Structure Form: Course, Semester, Amount, Due Date, Academic Year",
        "• Notice Form: Title, Content, Category (dropdown), Target Role, Publish toggle",
        "• Timetable Form: Course, Subject, Faculty, Day, Start Time, End Time, Room",
    ])

    _content_slide(prs, "6.1 Output Forms / Screens", [
        "• Admin Dashboard: Statistics cards (Students, Faculty, Courses, Notices count)",
        "• Student Dashboard: Attendance percentage with progress bar, Quick links",
        "• Faculty Dashboard: Subject assignments count, Unread messages, Quick actions",
        "• Student List: Searchable table with enrollment no, name, course, semester, actions",
        "• Attendance Report: Subject-wise summary with percentage charts (Chart.js)",
        "• Exam Results: Internal + External marks, Total, Grade — per student per subject",
        "• Fee Status: Paid/Pending indicators with amount, due date, receipt links",
        "• Timetable View: Day-wise grid with time slots, subjects, faculty, rooms",
        "• QR Session View: Live QR code display, countdown timer, scan count",
        "• Data Dictionary: Table schema viewer with export (PDF/DOCX/PPTX)",
    ])

    # ──────────────────────────────────────────────────────────
    # 7. IMPLEMENTATION (Slides 25-27)
    # ──────────────────────────────────────────────────────────
    _section_title_slide(prs, 7, "Implementation")

    _content_slide(prs, "7. Implementation — Architecture", [
        "• Architecture Pattern: MVC (Model-View-Controller)",
        "  – Model: PostgreSQL database with 20+ normalized tables",
        "  – View: Jinja2 HTML templates with Bootstrap 5 styling",
        "  – Controller: Flask route handlers organized in Blueprints",
        "",
        "• Backend Structure:",
        "  – app.py — Main Flask application and SocketIO setup",
        "  – config.py — Configuration management (DB, QR, Security)",
        "  – database.py — Database connection with context manager",
        "  – auth.py — Authentication and role-based authorization decorators",
        "  – routes/ — 12 Blueprint modules for each feature area",
        "  – qr_service.py — QR code generation and validation engine",
        "  – security_service.py — Fraud detection and audit logging",
    ])

    _content_slide(prs, "7. Implementation — Security", [
        "• Authentication:",
        "  – Werkzeug password hashing (PBKDF2-SHA256)",
        "  – Flask session management with configurable lifetime",
        "  – Login via username/email or enrollment number (students)",
        "",
        "• Authorization:",
        "  – @require_roles() decorator for role-based access control",
        "  – Four roles: Admin, Faculty, Student, Accountant",
        "",
        "• QR Security:",
        "  – Encrypted tokens using itsdangerous (time-limited signatures)",
        "  – Unique (session_id, student_id) constraint prevents duplicate scans",
        "  – IP address, user agent, and geolocation logging per scan",
        "  – Fraud detection dashboard with audit trail",
    ])

    _content_slide(prs, "7. Implementation — Deployment", [
        "• Development Environment:",
        "  – Python 3.10+ with pip and virtual environments",
        "  – MySQL 8.x for local development",
        "  – Flask development server with hot reload",
        "",
        "• Production Deployment (Render):",
        "  – PostgreSQL database (auto-detected via DATABASE_URL)",
        "  – Gunicorn WSGI server with Eventlet worker class",
        "  – Procfile: gunicorn --worker-class eventlet -w 1 -b 0.0.0.0:$PORT app:app",
        "",
        "• Environment Variables:",
        "  – FLASK_ENV, DATABASE_URL, SECRET_KEY",
        "  – QR_VALIDITY_MINUTES, GEOFENCE_RADIUS_KM",
    ])

    # ──────────────────────────────────────────────────────────
    # 8. BIBLIOGRAPHY (Slide 28)
    # ──────────────────────────────────────────────────────────
    _section_title_slide(prs, 8, "Bibliography")

    _content_slide(prs, "8. Bibliography — References", [
        "• Flask Official Documentation — https://flask.palletsprojects.com/",
        "• PostgreSQL Documentation — https://www.postgresql.org/docs/",
        "• Bootstrap 5 Documentation — https://getbootstrap.com/docs/5.3/",
        "• Chart.js Documentation — https://www.chartjs.org/docs/",
        "• python-pptx Documentation — https://python-pptx.readthedocs.io/",
        "• ReportLab User Guide — https://docs.reportlab.com/",
        "• Flask-SocketIO — https://flask-socketio.readthedocs.io/",
        "• QR Code Library — https://pypi.org/project/qrcode/",
        "• Werkzeug Security — https://werkzeug.palletsprojects.com/",
        "• Render Deployment Guide — https://render.com/docs/",
        "",
        "• Books & References:",
        "  – 'Flask Web Development' by Miguel Grinberg (O'Reilly)",
        "  – 'Database System Concepts' by Silberschatz, Korth, Sudarshan",
    ])

    # ──────────────────────────────────────────────────────────
    # THANK YOU SLIDE
    # ──────────────────────────────────────────────────────────
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, 11, 19, 43)
    _add_shape_rect(slide, Inches(0), Inches(2.5), Inches(13.333), Inches(2.8), ROYAL_BLUE)
    _add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(1.2),
                  "THANK YOU", font_size=54, color=WHITE, bold=True,
                  alignment=PP_ALIGN.CENTER)
    _add_shape_rect(slide, Inches(5.5), Inches(3.9), Inches(2.3), Inches(0.05), GOLD)
    _add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.8),
                  "College Management System — Questions & Discussion",
                  font_size=22, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    # Save to buffer
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
