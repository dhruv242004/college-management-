from flask import Blueprint, render_template, send_file, abort, flash, redirect, url_for, make_response, request
from auth import require_roles, get_current_user
from database import db_cursor
import io
import traceback
import xml.sax.saxutils as saxutils
from datetime import datetime

# Document Generation Libraries
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, PageBreak
from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from docx import Document
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

admin_bp = Blueprint('admin_bp', __name__)

@admin_bp.route('/data-dictionary')
@require_roles('admin')
def data_dictionary():
    user = get_current_user()
    
    with db_cursor() as (conn, cur):
        # Fetch all public base tables
        cur.execute("""
            SELECT table_name 
            FROM information_schema.tables 
            WHERE table_schema = 'public' 
            AND table_type = 'BASE TABLE'
            ORDER BY table_name;
        """)
        tables = cur.fetchall()
        
        data_dict = []
        for table in tables:
            tname = table['table_name']
            
            # Fetch column details for each table
            cur.execute("""
                SELECT 
                    column_name, 
                    data_type, 
                    is_nullable, 
                    column_default,
                    character_maximum_length
                FROM information_schema.columns 
                WHERE table_name = %s 
                AND table_schema = 'public'
                ORDER BY ordinal_position;
            """, (tname,))
            columns = cur.fetchall()
            
            data_dict.append({
                'table_name': tname,
                'columns': columns
            })
            
    return render_template('admin/data_dictionary.html', 
                          user=user, 
                          data_dict=data_dict)

@admin_bp.route('/pending-students')
@require_roles('admin')
def pending_students():
    """List students awaiting verification."""
    with db_cursor() as (conn, cur):
        cur.execute("""
            SELECT s.id, s.enrollment_no, s.first_name, s.last_name, s.email, 
                   c.name AS course_name, s.created_at
            FROM students s
            JOIN courses c ON c.id = s.course_id
            WHERE s.is_verified = FALSE
            ORDER BY s.created_at DESC
        """)
        pending = cur.fetchall()
    return render_template('admin/pending_students.html', pending=pending)

@admin_bp.route('/verify-student/<int:sid>', methods=['POST'])
@require_roles('admin')
def verify_student(sid):
    """Verify a student registration."""
    action = request.form.get('action')
    with db_cursor() as (conn, cur):
        if action == 'approve':
            cur.execute("UPDATE students SET is_verified = TRUE WHERE id = %s", (sid,))
            flash("Student account approved successfully.", "success")
        elif action == 'reject':
            # Optionally delete the user/student or just leave as is
            cur.execute("SELECT user_id FROM students WHERE id = %s", (sid,))
            row = cur.fetchone()
            if row:
                uid = row['user_id']
                cur.execute("DELETE FROM students WHERE id = %s", (sid,))
                cur.execute("DELETE FROM users WHERE id = %s", (uid,))
                flash("Student registration rejected and removed.", "info")
        conn.commit()
    return redirect(url_for('admin_bp.pending_students'))

@admin_bp.route('/generate-project-presentation')
@require_roles('admin')
def generate_project_presentation():
    """Generate a comprehensive project report in PPTX format."""
    print("DEBUG: Generating full project presentation")
    try:
        prs = PptxPresentation()
        
        # 1. Title Slide
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "College Management System"
        slide.placeholders[1].text = "Comprehensive Project Documentation & System Analysis\nCreated for Admin Review"

        # 2. Introduction
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "1. Introduction"
        content = slide.placeholders[1].text_frame
        content.text = "The College Management System is a digital platform designed to automate academic and administrative tasks."
        content.add_paragraph().text = "Key Goals:"
        p = content.add_paragraph()
        p.text = "- Streamline student enrollment and verification"
        p.level = 1
        p = content.add_paragraph()
        p.text = "- Modernize attendance tracking via QR codes"
        p.level = 1
        p = content.add_paragraph()
        p.text = "- Provide real-time data for admin decision making"
        p.level = 1

        # 3. System Architecture
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "2. System Architecture"
        content = slide.placeholders[1].text_frame
        content.text = "The system follows a Model-View-Controller (MVC) pattern using Flask."
        content.add_paragraph().text = "Frontend: HTML5, CSS3 (Glassmorphism), Bootstrap 5, JS"
        content.add_paragraph().text = "Backend: Python (Flask), Socket.IO for real-time chat"
        content.add_paragraph().text = "Database: PostgreSQL (Production) / MySQL (Local)"

        # 4. Database Schema (Summary)
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "3. Database Design (ER Summary)"
        content = slide.placeholders[1].text_frame
        content.text = "The system consists of several core modules:"
        content.add_paragraph().text = "- Users & Roles (RBAC)"
        content.add_paragraph().text = "- Students & Faculty Profiles"
        content.add_paragraph().text = "- Academic Records (Courses, Subjects, Exams)"
        content.add_paragraph().text = "- Financials (Fee Structures & Payments)"
        content.add_paragraph().text = "- Communication (Real-time Chat & Notices)"

        # 5. Data Dictionary (One Slide per major table)
        with db_cursor() as (conn, cur):
            cur.execute("""
                SELECT table_name 
                FROM information_schema.tables 
                WHERE table_schema = 'public' AND table_type = 'BASE TABLE'
                LIMIT 5;
            """)
            table_names = [row['table_name'] for row in cur.fetchall()]
            
            for tname in table_names:
                cur.execute("""
                    SELECT column_name, data_type, is_nullable
                    FROM information_schema.columns 
                    WHERE table_name = %s AND table_schema = 'public'
                    ORDER BY ordinal_position LIMIT 10;
                """, (tname,))
                cols = cur.fetchall()
                
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.title.text = f"Data Dictionary: {tname}"
                
                left, top, width, height = Inches(0.5), Inches(1.5), Inches(9.0), Inches(0.8)
                t_shape = slide.shapes.add_table(len(cols) + 1, 3, left, top, width, height)
                table = t_shape.table
                
                # Headers
                for i, h in enumerate(['Column', 'Type', 'Null']):
                    cell = table.cell(0, i)
                    cell.text = h
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(65, 105, 225)
                
                # Data
                for r, col in enumerate(cols):
                    table.cell(r+1, 0).text = str(col['column_name'])
                    table.cell(r+1, 1).text = str(col['data_type'])
                    table.cell(r+1, 2).text = str(col['is_nullable'])

        # 6. Conclusion
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Final Summary"
        content = slide.placeholders[1].text_frame
        content.text = "This presentation provides a technical overview of the College Management System."
        content.add_paragraph().text = "Generated automatically by the Admin Export Engine."

        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        
        return send_file(
            output,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name='College_Management_System_Report.pptx'
        )
            
    except Exception as e:
        print(f"CRITICAL: PPT generation failed: {str(e)}")
        traceback.print_exc()
        flash(f"Failed to generate presentation: {str(e)}", "error")
        return redirect(url_for('dashboard'))

@admin_bp.route('/export/<table_name>/<fmt>')
@require_roles('admin')
def export_table(table_name, fmt):
    print(f"DEBUG: Exporting table '{table_name}' in format '{fmt}'")
    try:
        with db_cursor() as (conn, cur):
            cur.execute("""
                SELECT column_name, data_type, is_nullable, column_default, character_maximum_length
                FROM information_schema.columns 
                WHERE table_name = %s AND table_schema = 'public'
                ORDER BY ordinal_position;
            """, (table_name,))
            columns = cur.fetchall()
            
        if not columns:
            print(f"DEBUG: No columns found for table '{table_name}'")
            flash(f"No metadata found for table: {table_name}", "warning")
            return redirect(url_for('admin_bp.data_dictionary'))

        output = io.BytesIO()
        safe_table_name = "".join([c if c.isalnum() else "_" for c in table_name])
        filename = f"{safe_table_name}_schema"
        mimetype = 'application/octet-stream'

        if fmt == 'pdf':
            mimetype = 'application/pdf'
            doc = SimpleDocTemplate(output, 
                                  pagesize=landscape(letter), 
                                  leftMargin=0.5*inch, 
                                  rightMargin=0.5*inch, 
                                  topMargin=0.5*inch, 
                                  bottomMargin=0.5*inch)
            elements = []
            styles = getSampleStyleSheet()
            cell_style = ParagraphStyle(name='CellStyle', parent=styles['Normal'], fontSize=8, leading=10, wordWrap='CJK')
            
            elements.append(Paragraph(f"Data Dictionary: {table_name.upper()}", styles['Title']))
            elements.append(Spacer(1, 12))
            
            headers = ['Column', 'Type', 'Null', 'Default', 'Max Len']
            data = [[Paragraph(f"<b>{h}</b>", cell_style) for h in headers]]
            
            for col in columns:
                data.append([
                    Paragraph(saxutils.escape(str(col['column_name'])), cell_style),
                    Paragraph(saxutils.escape(str(col['data_type'])), cell_style),
                    Paragraph(saxutils.escape(str(col['is_nullable'])), cell_style),
                    Paragraph(saxutils.escape(str(col['column_default']) if col['column_default'] else '-'), cell_style),
                    Paragraph(saxutils.escape(str(col['character_maximum_length']) if col['character_maximum_length'] else '-'), cell_style)
                ])
                
            col_widths = [1.5*inch, 1.5*inch, 0.8*inch, 5.0*inch, 1.2*inch]
            t = Table(data, colWidths=col_widths, hAlign='LEFT', repeatRows=1)
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4169E1')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.whitesmoke, colors.lightgrey]),
            ]))
            elements.append(t)
            doc.build(elements)

        elif fmt == 'docx':
            mimetype = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            doc = Document()
            doc.add_heading(f'Data Dictionary: {table_name}', 0)
            table = doc.add_table(rows=1, cols=5)
            table.style = 'Table Grid'
            hdr_cells = table.rows[0].cells
            for i, h in enumerate(['Column Name', 'Data Type', 'Nullable', 'Default', 'Max Length']):
                hdr_cells[i].text = h
            
            for col in columns:
                row_cells = table.add_row().cells
                row_cells[0].text = str(col['column_name'])
                row_cells[1].text = str(col['data_type'])
                row_cells[2].text = str(col['is_nullable'])
                row_cells[3].text = str(col['column_default']) if col['column_default'] else '-'
                row_cells[4].text = str(col['character_maximum_length']) if col['character_maximum_length'] else '-'
            doc.save(output)

        elif fmt == 'pptx':
            mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            prs = PptxPresentation()
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = f"Data Dictionary: {table_name}"
            left, top, width, height = Inches(0.5), Inches(1.5), Inches(9.0), Inches(0.8)
            table_shape = slide.shapes.add_table(len(columns) + 1, 5, left, top, width, height)
            table = table_shape.table
            
            headers = ['Column', 'Type', 'Null', 'Default', 'Max Len']
            for i, h in enumerate(headers):
                cell = table.cell(0, i)
                cell.text = h
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(65, 105, 225)
            
            for r, col in enumerate(columns):
                table.cell(r+1, 0).text = str(col['column_name'])
                table.cell(r+1, 1).text = str(col['data_type'])
                table.cell(r+1, 2).text = str(col['is_nullable'])
                table.cell(r+1, 3).text = str(col['column_default']) if col['column_default'] else '-'
                table.cell(r+1, 4).text = str(col['character_maximum_length']) if col['character_maximum_length'] else '-'
            prs.save(output)

        else:
            abort(400)

        # Build Hardened Response
        output.seek(0)
        response = make_response(output.getvalue())
        response.headers['Content-Type'] = mimetype
        response.headers['Content-Disposition'] = f'attachment; filename="{filename}.{fmt}"'
        print(f"DEBUG: Successfully generated {fmt} for {table_name}")
        return response
            
    except Exception as e:
        print(f"CRITICAL: Export failed for {table_name} ({fmt}): {str(e)}")
        traceback.print_exc()
        flash(f"System error during export: {str(e)}", "error")
        return redirect(url_for('admin_bp.data_dictionary'))

@admin_bp.route('/export-all/<fmt>')
@require_roles('admin')
def export_all_tables(fmt):
    print(f"DEBUG: Exporting ALL tables in format '{fmt}'")
    try:
        with db_cursor() as (conn, cur):
            cur.execute("""
                SELECT table_name 
                FROM information_schema.tables 
                WHERE table_schema = 'public' AND table_type = 'BASE TABLE'
                ORDER BY table_name;
            """)
            table_names = [row['table_name'] for row in cur.fetchall()]
            
            all_metadata = []
            for tname in table_names:
                cur.execute("""
                    SELECT column_name, data_type, is_nullable, column_default, character_maximum_length
                    FROM information_schema.columns 
                    WHERE table_name = %s AND table_schema = 'public'
                    ORDER BY ordinal_position;
                """, (tname,))
                all_metadata.append({'table_name': tname, 'columns': cur.fetchall()})

        if not all_metadata:
            flash("No tables found to export.", "warning")
            return redirect(url_for('admin_bp.data_dictionary'))

        output = io.BytesIO()
        filename = "full_database_dictionary"
        mimetype = 'application/octet-stream'
        
        if fmt == 'pdf':
            mimetype = 'application/pdf'
            doc = SimpleDocTemplate(output, 
                                  pagesize=landscape(letter), 
                                  leftMargin=0.5*inch, 
                                  rightMargin=0.5*inch, 
                                  topMargin=0.5*inch, 
                                  bottomMargin=0.5*inch)
            elements = []
            styles = getSampleStyleSheet()
            cell_style = ParagraphStyle(name='CellStyle', parent=styles['Normal'], fontSize=8, leading=10, wordWrap='CJK')
            
            elements.append(Paragraph("System Data Dictionary - Full Report", styles['Title']))
            elements.append(Paragraph(f"Generated on: {datetime.now().strftime('%Y-%m-%d %H:%M %p')}", styles['Normal']))
            elements.append(Spacer(1, 24))
            
            for i, table in enumerate(all_metadata):
                if i > 0: elements.append(PageBreak())
                elements.append(Paragraph(f"Table: {table['table_name'].upper()}", styles['Heading2']))
                elements.append(Spacer(1, 10))
                
                headers = ['Column', 'Type', 'Null', 'Default', 'Max Len']
                data = [[Paragraph(f"<b>{h}</b>", cell_style) for h in headers]]
                for col in table['columns']:
                    data.append([
                        Paragraph(saxutils.escape(str(col['column_name'])), cell_style),
                        Paragraph(saxutils.escape(str(col['data_type'])), cell_style),
                        Paragraph(saxutils.escape(str(col['is_nullable'])), cell_style),
                        Paragraph(saxutils.escape(str(col['column_default']) if col['column_default'] else '-'), cell_style),
                        Paragraph(saxutils.escape(str(col['character_maximum_length']) if col['character_maximum_length'] else '-'), cell_style)
                    ])
                
                t = Table(data, colWidths=[1.5*inch, 1.5*inch, 0.8*inch, 5.0*inch, 1.2*inch], repeatRows=1)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#4169E1')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.whitesmoke, colors.lightgrey]),
                ]))
                elements.append(t)
                elements.append(Spacer(1, 20))
            doc.build(elements)

        elif fmt == 'docx':
            mimetype = 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
            doc = Document()
            doc.add_heading('System Data Dictionary - Full Report', 0)
            for i, table in enumerate(all_metadata):
                if i > 0: doc.add_page_break()
                doc.add_heading(f"Table: {table['table_name']}", level=1)
                t = doc.add_table(rows=1, cols=5)
                t.style = 'Table Grid'
                hdr_cells = t.rows[0].cells
                for idx, h in enumerate(['Column', 'Type', 'Null', 'Default', 'Max Len']):
                    hdr_cells[idx].text = h
                for col in table['columns']:
                    row_cells = t.add_row().cells
                    row_cells[0].text = str(col['column_name'])
                    row_cells[1].text = str(col['data_type'])
                    row_cells[2].text = str(col['is_nullable'])
                    row_cells[3].text = str(col['column_default']) if col['column_default'] else '-'
                    row_cells[4].text = str(col['character_maximum_length']) if col['character_maximum_length'] else '-'
            doc.save(output)

        elif fmt == 'pptx':
            mimetype = 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
            prs = PptxPresentation()
            for table in all_metadata:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.shapes.title.text = f"Schema: {table['table_name']}"
                left, top, width, height = Inches(0.5), Inches(1.5), Inches(9.0), Inches(0.8)
                t_shape = slide.shapes.add_table(len(table['columns']) + 1, 5, left, top, width, height)
                t = t_shape.table
                for idx, h in enumerate(['Column', 'Type', 'Null', 'Default', 'Max Len']):
                    cell = t.cell(0, idx)
                    cell.text = h
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(65, 105, 225)
                for r, col in enumerate(table['columns']):
                    t.cell(r+1, 0).text = str(col['column_name'])
                    t.cell(r+1, 1).text = str(col['data_type'])
                    t.cell(r+1, 2).text = str(col['is_nullable'])
                    t.cell(r+1, 3).text = str(col['column_default']) if col['column_default'] else '-'
                    t.cell(r+1, 4).text = str(col['character_maximum_length']) if col['character_maximum_length'] else '-'
            prs.save(output)

        else:
            abort(400)

        # Build Hardened Response
        output.seek(0)
        response = make_response(output.getvalue())
        response.headers['Content-Type'] = mimetype
        response.headers['Content-Disposition'] = f'attachment; filename="{filename}.{fmt}"'
        print(f"DEBUG: Successfully generated Bulk {fmt}")
        return response
            
    except Exception as e:
        print(f"CRITICAL: Bulk export failed ({fmt}): {str(e)}")
        traceback.print_exc()
        flash(f"Bulk export failed: {str(e)}", "error")
        return redirect(url_for('admin_bp.data_dictionary'))
